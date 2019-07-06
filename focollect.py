import subprocess
import sys

def install(package):
    subprocess.call([sys.executable, "-m", "pip", "install", package])

try:
	from pptx import Presentation
	from pptx.util import Cm, Pt

except:
    install('python-pptx')

from pptx import Presentation
from pptx.util import Cm, Pt
import os

prs = Presentation()

pic_left = Cm(2.37)
pic_top = Cm(4.52)
pic_width = Cm(21.01)
pic_height = Cm(11.81)

FileList = os.listdir()
print(FileList)
for File in FileList:
	if File[-4:] == ".jpg" or File[-4:] == ".png":
		print(File)	
		#Делим материал для заголовка и подзаголовка по $$$
		ListOfNames = File[:-4].split("$$$")
		title_slide_layout = prs.slide_layouts[0]
		slide = prs.slides.add_slide(title_slide_layout)	

		#Выбираем title и задаем координаты
		title = slide.shapes.title
		title.text = ListOfNames[0]
		title.left = Cm(1.71)
		title.top = Cm(0.54)
		title.width = Cm(23.26)
		title.height = Cm(2.26)	

		#Выбираем subtitle и задаем координаты
		subtitle = slide.placeholders[1]
		subtitle.left = Cm(3.99)
		subtitle.top = Cm(2.8)
		subtitle.width = Cm(17.78)
		subtitle.height = Cm(4.87)

		try:
			subtitle.text = ListOfNames[1]
		except:
			#Если вдруг слайды названы без $$$, мы прячем подзаголовок за картинку
			subtitle.text = ""
			subtitle.top = Cm(5)



		slide.shapes.add_picture(File, pic_left , pic_top, pic_width, pic_height)

prs.save('ФО.pptx')